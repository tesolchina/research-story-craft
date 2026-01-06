#!/usr/bin/env python3
"""
Script to convert all course materials to markdown format and maintain directory structure
"""
import os
import sys
import subprocess
import shutil
from pathlib import Path
import time
from datetime import datetime

# Define source and target directories
SOURCE_DIR = Path("/Users/simonwang/Documents/Usage/gtd001/Teaching/courses/MCCP 6020/materials/MCCP6020_2025-26sem1_TEACHER'S Folder")
TARGET_DIR = Path("/Users/simonwang/Documents/Usage/gtd001/Teaching/courses/MCCP 6020/materials/MachineReadableMat")
LOG_FILE = TARGET_DIR / "process.log"

def log_message(message):
    """Write message to log file and print to console"""
    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    log_message = f"[{timestamp}] {message}\n"
    print(log_message.strip())
    
    with open(LOG_FILE, 'a', encoding='utf-8') as f:
        f.write(log_message)

def ensure_target_structure(source_file_path):
    """Ensure the directory structure exists in the target directory"""
    # Get the relative path from source directory
    rel_path = source_file_path.relative_to(SOURCE_DIR)
    # Create the target path in MachineReadableMat
    target_path = TARGET_DIR / rel_path.parent
    target_path.mkdir(parents=True, exist_ok=True)
    return target_path

def convert_docx_to_md(docx_path, target_dir):
    """Convert DOCX file to markdown using pandoc"""
    try:
        # Get the base name without extension
        base_name = docx_path.stem
        # Create target markdown file path
        md_path = target_dir / f"{base_name}.md"
        
        # Run pandoc to convert docx to markdown
        cmd = ["pandoc", str(docx_path), "-f", "docx", "-t", "markdown", "-o", str(md_path)]
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        
        if result.returncode == 0:
            log_message(f"✓ Converted {docx_path.name} to {md_path.name}")
            return str(md_path)
        else:
            log_message(f"✗ Error converting {docx_path.name}: {result.stderr}")
            return None
    except subprocess.TimeoutExpired:
        log_message(f"✗ Timeout converting {docx_path.name}")
        return None
    except Exception as e:
        log_message(f"✗ Error converting {docx_path.name}: {str(e)}")
        return None

def convert_pdf_to_md(pdf_path, target_dir):
    """Convert PDF file to markdown using existing pdf_extract_md.py script"""
    try:
        # Get the base name without extension
        base_name = pdf_path.stem
        # Create target markdown file path
        md_path = target_dir / f"{base_name}.md"
        
        # Run the existing pdf extraction script
        cmd = ["python3", "/Users/simonwang/Documents/Usage/gtd001/tools/pdfWordConvert/pdf_extract_md.py", str(pdf_path), str(md_path)]
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        
        if result.returncode == 0:
            log_message(f"✓ Converted {pdf_path.name} to {md_path.name}")
            return str(md_path)
        else:
            log_message(f"✗ Error converting {pdf_path.name}: {result.stderr}")
            return None
    except subprocess.TimeoutExpired:
        log_message(f"✗ Timeout converting {pdf_path.name}")
        return None
    except Exception as e:
        log_message(f"✗ Error converting {pdf_path.name}: {str(e)}")
        return None

def convert_pptx_to_md(pptx_path, target_dir):
    """Convert PPTX file to markdown by first converting to PDF then to markdown"""
    try:
        # First convert PPTX to PDF using LibreOffice (if available)
        # This is a fallback approach since direct PPTX to Markdown conversion with pandoc isn't reliable
        base_name = pptx_path.stem
        pdf_path = target_dir / f"{base_name}_temp.pdf"
        
        # Try to convert PPTX to PDF using LibreOffice if available
        libreoffice_cmd = ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", str(target_dir), str(pptx_path)]
        try:
            result = subprocess.run(libreoffice_cmd, capture_output=True, text=True, timeout=180)
            if result.returncode == 0:
                # Find the converted PDF (it may have a slightly different name)
                converted_pdf = None
                for file in target_dir.iterdir():
                    if file.name.startswith(base_name) and file.suffix.lower() == '.pdf':
                        converted_pdf = file
                        break
                
                if converted_pdf:
                    # Now convert the PDF to markdown
                    md_path = target_dir / f"{base_name}.md"
                    # Rename the converted PDF to our target name temporarily
                    converted_pdf.rename(pdf_path)
                    md_result = convert_pdf_to_md(pdf_path, target_dir)
                    # Remove the temporary PDF after conversion
                    if pdf_path.exists():
                        pdf_path.unlink()
                    return md_result
                else:
                    log_message(f"✗ Could not find converted PDF for {pptx_path.name}")
                    return None
            else:
                log_message(f"✗ Error converting {pptx_path.name} to PDF: {result.stderr}")
                return None
        except FileNotFoundError:
            log_message(f"✗ LibreOffice not found, trying alternative PPTX conversion for {pptx_path.name}")
            # Try alternative method using python-pptx if available
            return convert_pptx_to_text_alternative(pptx_path, target_dir)
    except Exception as e:
        log_message(f"✗ Error converting {pptx_path.name}: {str(e)}")
        return None


def convert_pptx_to_text_alternative(pptx_path, target_dir):
    """Alternative method to convert PPTX to text/markdown using python-pptx"""
    try:
        from pptx import Presentation
        
        prs = Presentation(pptx_path)
        base_name = pptx_path.stem
        md_path = target_dir / f"{base_name}.md"
        
        slides_content = []
        for i, slide in enumerate(prs.slides, 1):
            slide_content = [f"## Slide {i}\n"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_content.append(shape.text)
            
            # Extract content from tables if present
            for shape in slide.shapes:
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = "| " + " | ".join([cell.text for cell in row.cells]) + " |"
                        slide_content.append(row_text)
            
            slides_content.append("\n".join(slide_content))
        
        # Write to markdown file
        with open(md_path, 'w', encoding='utf-8') as f:
            f.write("\n\n---\n\n".join(slides_content))
        
        log_message(f"✓ Converted {pptx_path.name} to {md_path.name} using alternative method")
        return str(md_path)
    except ImportError:
        log_message(f"✗ python-pptx not available, skipping PPTX conversion for {pptx_path.name}")
        return None
    except Exception as e:
        log_message(f"✗ Error in alternative PPTX conversion for {pptx_path.name}: {str(e)}")
        return None


def should_convert_file(file_path):
    """Check if the file should be converted based on its extension"""
    ext = file_path.suffix.lower()
    return ext in ['.docx', '.pdf', '.pptx']

def process_directory(source_dir):
    """Recursively process all files in the source directory"""
    total_files = 0
    converted_files = 0
    
    # First, count all files to convert
    for root, dirs, files in os.walk(source_dir):
        for file in files:
            file_path = Path(root) / file
            if should_convert_file(file_path):
                total_files += 1
    
    log_message(f"Found {total_files} files to convert")
    
    # Process each file
    for root, dirs, files in os.walk(source_dir):
        for file in files:
            file_path = Path(root) / file
            if should_convert_file(file_path):
                # Ensure target directory structure exists
                target_dir = ensure_target_structure(file_path)
                
                log_message(f"Processing: {file_path.name}")
                
                # Convert based on file extension
                result = None
                if file_path.suffix.lower() == '.docx':
                    result = convert_docx_to_md(file_path, target_dir)
                elif file_path.suffix.lower() == '.pdf':
                    result = convert_pdf_to_md(file_path, target_dir)
                elif file_path.suffix.lower() == '.pptx':
                    result = convert_pptx_to_md(file_path, target_dir)
                
                if result:
                    converted_files += 1
                    
                # Small delay to avoid overwhelming the system
                time.sleep(0.1)
    
    log_message(f"Conversion complete. {converted_files}/{total_files} files converted.")
    return converted_files, total_files

def main():
    log_message("Starting conversion of MCCP 6020 course materials to markdown...")
    log_message(f"Source directory: {SOURCE_DIR}")
    log_message(f"Target directory: {TARGET_DIR}")
    
    # Check if source directory exists
    if not SOURCE_DIR.exists():
        log_message(f"ERROR: Source directory does not exist: {SOURCE_DIR}")
        return 1
    
    # Create target directory if it doesn't exist
    TARGET_DIR.mkdir(parents=True, exist_ok=True)
    
    # Process all files
    converted, total = process_directory(SOURCE_DIR)
    
    log_message(f"Conversion process finished. Converted {converted} out of {total} files.")
    return 0

if __name__ == "__main__":
    sys.exit(main())